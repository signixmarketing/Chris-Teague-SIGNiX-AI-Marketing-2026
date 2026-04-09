#!/usr/bin/env python3
"""Build SIGNIX-PAID-MEDIA-SLIDES-CEO.pptx — trimmed for executive review.

Differences from the full deck:
  Slide 4 — keyword detail removed; shows bucket names + allocations only
  Slide 5 — 14 indicators reduced to 3 headline KPIs + 60-day check-in commitment
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colors ────────────────────────────────────────────────────────────────────
GREEN  = RGBColor(0x6d, 0xa3, 0x4a)
INK    = RGBColor(0x2e, 0x34, 0x40)
BODY   = RGBColor(0x54, 0x54, 0x54)
MUTED  = RGBColor(0x6b, 0x72, 0x80)
WHITE  = RGBColor(0xff, 0xff, 0xff)
CANVAS = RGBColor(0xf8, 0xfa, 0xfb)
RULE   = RGBColor(0xd8, 0xde, 0xe9)
LIGHT_BODY = RGBColor(0xc8, 0xd0, 0xd8)

FONT = "Calibri"
W    = Inches(10)
H    = Inches(5.625)

# ── Helpers ───────────────────────────────────────────────────────────────────

def prs_new():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color, border_color=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if border_color:
        s.line.color.rgb = border_color
    else:
        s.line.fill.background()
    return s

def txb(slide, l, t, w, h, text, size=11, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    color = color or BODY
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name   = FONT
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def mixed_txb(slide, l, t, w, h, segments, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    for text, size, bold, color, italic in segments:
        run = p.add_run()
        run.text = text
        run.font.name   = FONT
        run.font.size   = Pt(size)
        run.font.bold   = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box

def bullet_card(slide, l, t, w, h, header, items,
                bg_color=CANVAS, text_color=BODY):
    rect(slide, l, t, w, h, bg_color)
    rect(slide, l, t, Inches(0.045), h, GREEN)
    box = slide.shapes.add_textbox(l + Inches(0.12), t + Inches(0.14),
                                   w - Inches(0.2), h - Inches(0.22))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(4)
    run = p.add_run()
    run.text = header.upper()
    run.font.name  = FONT
    run.font.size  = Pt(8.5)
    run.font.bold  = True
    run.font.color.rgb = GREEN
    for item in items:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = "›  " + item
        run.font.name  = FONT
        run.font.size  = Pt(10)
        run.font.color.rgb = text_color

def top_bar(slide):
    rect(slide, Inches(0), Inches(0), W, Inches(0.055), GREEN)

def header(slide, eyebrow, title, title_size=20):
    txb(slide, Inches(0.65), Inches(0.09), Inches(8.7), Inches(0.22),
        eyebrow.upper(), size=7.5, bold=True, color=GREEN)
    txb(slide, Inches(0.65), Inches(0.31), Inches(8.7), Inches(0.6),
        title, size=title_size, bold=True, color=INK)

def footer(slide):
    fy = H - Inches(0.46)
    rect(slide, Inches(0), fy, W, Inches(0.46), WHITE)
    rect(slide, Inches(0), fy, W, Pt(0.75), RULE)
    txb(slide, W - Inches(1.55), fy + Inches(0.1),
        Inches(0.9), Inches(0.28),
        "SIGNiX", size=12, bold=True, color=INK, align=PP_ALIGN.RIGHT)
    rect(slide, W - Inches(0.6), fy + Inches(0.14),
         Inches(0.055), Inches(0.055), GREEN)


# ── SLIDE 1 — Title ───────────────────────────────────────────────────────────
prs = prs_new()
s1  = blank(prs)
bg(s1, INK)
rect(s1, Inches(0), H - Inches(0.055), W, Inches(0.055), GREEN)
rect(s1, Inches(3.82), Inches(0.9), Inches(0.07), Inches(0.68), GREEN)
txb(s1, Inches(3.95), Inches(0.86), Inches(4.0), Inches(0.78),
    "SIGNiX", size=42, bold=True, color=WHITE)
rect(s1, Inches(3.82), Inches(1.66), Inches(4.0), Inches(0.022), GREEN)
txb(s1, Inches(1.5), Inches(1.85), Inches(7.0), Inches(0.85),
    "2026 Paid Media Strategy",
    size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s1, Inches(1.5), Inches(2.78), Inches(7.0), Inches(0.44),
    "Moving Budget to the Deals That Matter",
    size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
txb(s1, Inches(1.5), Inches(3.3), Inches(7.0), Inches(0.32),
    "Chris  ·  Head of Growth and Marketing  ·  April 2026",
    size=10, color=MUTED, align=PP_ALIGN.CENTER)


# ── SLIDE 2 — The Problem ─────────────────────────────────────────────────────
s2 = blank(prs)
bg(s2, WHITE)
top_bar(s2)
header(s2, "Strategic Context", "Notary Leads Cost More Than They Bring In")
footer(s2)

bullet_card(s2,
    l=Inches(0.42), t=Inches(1.0), w=Inches(4.52), h=Inches(4.0),
    header="Current State: The Problem",
    items=[
        "$5K/month goes to individual notary searches on Google",
        "Low value per customer. Most pay once and don't come back.",
        "Each deal takes a sales rep and Gina's full onboarding time. Cost likely exceeds revenue.",
        "No LinkedIn presence. B2B buyers who aren't searching are out of reach.",
        "RON leads are the whole pipeline. That won't scale.",
    ],
    bg_color=INK,
    text_color=LIGHT_BODY,
)
bullet_card(s2,
    l=Inches(5.06), t=Inches(1.0), w=Inches(4.52), h=Inches(4.0),
    header="Proposed State: The Shift",
    items=[
        "Same $5K/month. No budget increase needed.",
        "Reach FI ops, IT, and compliance buyers. Higher value and recurring revenue.",
        "Reach ISV developers through the Flex API. One integration can multiply your transaction base.",
        "Add LinkedIn to reach B2B buyers before they start searching.",
        "Keep RON only for institutional buyers.",
    ],
    bg_color=CANVAS,
    text_color=BODY,
)


# ── SLIDE 3 — Budget ──────────────────────────────────────────────────────────
s3 = blank(prs)
bg(s3, WHITE)
top_bar(s3)
header(s3, "Budget", "$5,000/Month, Moved to Enterprise Buyers")
footer(s3)

stat_data = [
    ("$3,000", "GOOGLE SEARCH ADS",
     "Five keyword groups targeting FI buyers and ISV developers. "
     "Focused on competitor displacement, authentication, and regulated industry terms."),
    ("$1,500", "LINKEDIN SPONSORED",
     "Two groups: FI ops and compliance leaders ($800) and ISV product teams "
     "($700). Sponsored posts only for the first two months."),
    ("$500",   "TEST & RESERVE",
     "Used for landing page tests, new keyword trials, and cost-per-click "
     "buffer. Moved to the top channel after a 30-day review."),
]
bw  = Inches(2.8)
gap = Inches(0.35)
bl  = Inches(0.65)

for i, (amount, label, desc) in enumerate(stat_data):
    x  = bl + i * (bw + gap)
    by = Inches(1.08)
    bh = Inches(2.85)
    rect(s3, x, by, bw, bh, CANVAS)
    rect(s3, x, by + bh - Inches(0.04), bw, Inches(0.04), GREEN)
    txb(s3, x, by + Inches(0.15), bw, Inches(0.6),
        amount, size=28, bold=True, color=INK, align=PP_ALIGN.CENTER)
    txb(s3, x, by + Inches(0.75), bw, Inches(0.25),
        label, size=7.5, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    txb(s3, x + Inches(0.12), by + Inches(1.05), bw - Inches(0.24), Inches(1.65),
        desc, size=10, color=BODY)

ny = Inches(4.05)
rect(s3, Inches(0.65), ny, Inches(8.7), Inches(0.76), CANVAS)
rect(s3, Inches(0.65), ny, Inches(0.045), Inches(0.76), GREEN)
mixed_txb(s3, Inches(0.75), ny + Inches(0.12), Inches(8.45), Inches(0.55),
    segments=[
        ("No budget increase needed.  ", 10.5, True,  INK,  False),
        ("We are moving existing spend, not adding to it. "
         "Notary keywords stop right away. RON stays at a small budget "
         "for institutional buyers only.", 10.5, False, BODY, False),
    ]
)


# ── SLIDE 4 — Channels (CEO trimmed) ─────────────────────────────────────────
# Keyword table shows bucket + allocation only (no keyword detail)
# LinkedIn cards kept in full
s4 = blank(prs)
bg(s4, WHITE)
top_bar(s4)
header(s4, "Channels & Keywords", "Google + LinkedIn: Reaching the Right Buyers")
footer(s4)

# Simplified bucket table — 2 columns only
buckets = [
    ("Competitor Displacement",  "~$900"),
    ("Industry / Use-Case",      "~$1,000"),
    ("Authentication / Security","~$600"),
    ("Flex API / Developer",     "~$500"),
    ("RON — Institutional Only", "~$300"),
]
col_w = [Inches(2.6), Inches(0.8)]
tl = Inches(0.42)
tt = Inches(1.02)
tw = sum(col_w)
th = Inches(4.0)

tbl_shape = s4.shapes.add_table(len(buckets) + 1, 2, tl, tt, tw, th)
tbl = tbl_shape.table

for ci, h in enumerate(["Google Keyword Focus", "Monthly Budget"]):
    cell = tbl.cell(0, ci)
    cell.fill.solid()
    cell.fill.fore_color.rgb = INK
    cell.text = h
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name  = FONT
    run.font.size  = Pt(9)
    run.font.bold  = True
    run.font.color.rgb = WHITE

for ri, (bucket, alloc) in enumerate(buckets):
    row_bg = CANVAS if ri % 2 == 0 else WHITE
    for ci, val in enumerate([bucket, alloc]):
        cell = tbl.cell(ri + 1, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = row_bg
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name  = FONT
        run.font.size  = Pt(10)
        run.font.bold  = (ci == 0)
        run.font.color.rgb = INK if ci == 0 else GREEN

for ri in range(len(buckets) + 1):
    for ci, cw in enumerate(col_w):
        tbl.cell(ri, ci).width = cw

# LinkedIn cards
cx = Inches(4.1)
bullet_card(s4,
    l=cx, t=Inches(1.02), w=Inches(5.46), h=Inches(1.9),
    header="LinkedIn — FI Audience ($800)",
    items=[
        "VP Operations, IT Director, Director of Lending, CCO",
        "Banks and credit unions — 50–5,000 employees",
        "Show them how SIGNiX stops fraud and proves compliance",
    ],
    bg_color=CANVAS, text_color=BODY,
)
bullet_card(s4,
    l=cx, t=Inches(3.06), w=Inches(5.46), h=Inches(1.9),
    header="LinkedIn — ISV / Flex API ($700)",
    items=[
        "Product Manager, VP of Product, CTO, VP Engineering",
        "Fintech, LOS platforms, banking software — 10–500 employees",
        "Show them a white-label API built for regulated industries",
    ],
    bg_color=CANVAS, text_color=BODY,
)


# ── SLIDE 5 — Measurement (CEO trimmed) ──────────────────────────────────────
# 3 headline KPIs + 60-day check-in commitment. No 14-item indicator list.
s5 = blank(prs)
bg(s5, WHITE)
top_bar(s5)
header(s5, "Measurement", "Three Numbers That Tell Us If It's Working")
footer(s5)

kpis = [
    ("Cost Per\nQualified Lead",
     "Compare new B2B leads against what we pay per notary lead today. "
     "If it's lower — and we expect it will be — the case is proven."),
    ("Staff Hours\nPer Closed Deal",
     "The hidden cost metric. Notary deals consume rep and CSM time well "
     "beyond their revenue. B2B deals should close at lower internal cost."),
    ("Pipeline Sourced\nin HubSpot",
     "Total revenue potential from paid-sourced leads, tracked by channel. "
     "This is the number the CEO sees at 90 days."),
]

kw  = Inches(2.7)
kgap = Inches(0.45)
kl  = Inches(0.65)
ky  = Inches(1.1)
kh  = Inches(3.2)

for i, (title, desc) in enumerate(kpis):
    x = kl + i * (kw + kgap)
    rect(s5, x, ky, kw, kh, CANVAS)
    rect(s5, x, ky + kh - Inches(0.05), kw, Inches(0.05), GREEN)
    # Number badge
    rect(s5, x + kw/2 - Inches(0.28), ky + Inches(0.18),
         Inches(0.56), Inches(0.56), GREEN)
    txb(s5, x + kw/2 - Inches(0.28), ky + Inches(0.18),
        Inches(0.56), Inches(0.56),
        str(i + 1), size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Title
    txb(s5, x + Inches(0.15), ky + Inches(0.85),
        kw - Inches(0.3), Inches(0.65),
        title, size=13, bold=True, color=INK, align=PP_ALIGN.CENTER)
    # Description
    txb(s5, x + Inches(0.15), ky + Inches(1.6),
        kw - Inches(0.3), Inches(1.45),
        desc, size=10, color=BODY)

# 60-day commitment note
ny = Inches(4.42)
rect(s5, Inches(0.65), ny, Inches(8.7), Inches(0.72), INK)
rect(s5, Inches(0.65), ny, Inches(0.045), Inches(0.72), GREEN)
txb(s5, Inches(0.76), ny + Inches(0.13), Inches(8.4), Inches(0.5),
    "60-day check-in: cost per qualified lead vs. notary lead baseline, "
    "staff hours per deal, and pipeline sourced. If the numbers don't hold up, we adjust.",
    size=10.5, color=WHITE)


# ── Save ──────────────────────────────────────────────────────────────────────
out = ("/Users/chris/Desktop/AI Summit 2026 Q2/"
       "proj-template-and-lease-SIGNiX-app/PROJECT-DOCS/"
       "SIGNIX-PAID-MEDIA-SLIDES-CEO.pptx")
prs.save(out)
print(f"Saved: {out}")
