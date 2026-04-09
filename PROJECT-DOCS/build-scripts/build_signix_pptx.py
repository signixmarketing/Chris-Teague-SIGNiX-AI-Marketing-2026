#!/usr/bin/env python3
"""Build SIGNIX-PAID-MEDIA-SLIDES.pptx with SIGNiX branding."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colors ────────────────────────────────────────────────────────────────────
GREEN  = RGBColor(0x6d, 0xa3, 0x4a)
INK    = RGBColor(0x2e, 0x34, 0x40)
BODY   = RGBColor(0x54, 0x54, 0x54)
MUTED  = RGBColor(0x6b, 0x72, 0x80)
WHITE  = RGBColor(0xff, 0xff, 0xff)
CANVAS = RGBColor(0xf8, 0xfa, 0xfb)
PANEL  = RGBColor(0xec, 0xef, 0xf4)
RULE   = RGBColor(0xd8, 0xde, 0xe9)
LIGHT_BODY = RGBColor(0xc8, 0xd0, 0xd8)

FONT = "Calibri"
W    = Inches(10)
H    = Inches(5.625)

# ── Core helpers ──────────────────────────────────────────────────────────────

def prs_new():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color, border_color=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if border_color:
        s.line.color.rgb = border_color
    else:
        s.line.fill.background()
    return s

def txb(slide, l, t, w, h, text, size=11, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    color = color or BODY
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = FONT
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def mixed_txb(slide, l, t, w, h, segments, wrap=True):
    """Text box with multiple runs. segments = [(text, size, bold, color, italic)]"""
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    for text, size, bold, color, italic in segments:
        run = p.add_run()
        run.text = text
        run.font.name  = FONT
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box

def bullet_card(slide, l, t, w, h, header, items,
                bg_color=CANVAS, text_color=BODY):
    rect(slide, l, t, w, h, bg_color)
    rect(slide, l, t, Inches(0.045), h, GREEN)

    box = slide.shapes.add_textbox(l + Inches(0.12), t + Inches(0.14),
                                   w - Inches(0.2), h - Inches(0.22))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True

    # Header
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(4)
    run = p.add_run()
    run.text = header.upper()
    run.font.name  = FONT
    run.font.size  = Pt(8.5)
    run.font.bold  = True
    run.font.color.rgb = GREEN

    # Bullets
    for item in items:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = "›  " + item
        run.font.name  = FONT
        run.font.size  = Pt(10)
        run.font.color.rgb = text_color

def top_bar(slide):
    rect(slide, Inches(0), Inches(0), W, Inches(0.055), GREEN)

def header(slide, eyebrow, title, title_size=20):
    txb(slide, Inches(0.65), Inches(0.09), Inches(8.7), Inches(0.22),
        eyebrow.upper(), size=7.5, bold=True, color=GREEN)
    txb(slide, Inches(0.65), Inches(0.31), Inches(8.7), Inches(0.6),
        title, size=title_size, bold=True, color=INK)

def footer(slide):
    fy = H - Inches(0.46)
    rect(slide, Inches(0), fy, W, Inches(0.46), WHITE)
    rect(slide, Inches(0), fy, W, Pt(0.75), RULE)
    # SIGNiX text logo
    txb(slide, W - Inches(1.55), fy + Inches(0.1),
        Inches(0.9), Inches(0.28),
        "SIGNiX", size=12, bold=True, color=INK, align=PP_ALIGN.RIGHT)
    # green i-dot accent
    rect(slide, W - Inches(0.6), fy + Inches(0.14),
         Inches(0.055), Inches(0.055), GREEN)


# ── SLIDE 1 — Title ───────────────────────────────────────────────────────────
prs = prs_new()
s1  = blank(prs)
bg(s1, INK)

# Green bottom rule
rect(s1, Inches(0), H - Inches(0.055), W, Inches(0.055), GREEN)

# Logo area — centered text mark
rect(s1, Inches(3.82), Inches(0.9), Inches(0.07), Inches(0.68), GREEN)
txb(s1, Inches(3.95), Inches(0.86), Inches(4.0), Inches(0.78),
    "SIGNiX", size=42, bold=True, color=WHITE)

# Thin green rule under logo
rect(s1, Inches(3.82), Inches(1.66), Inches(4.0), Inches(0.022), GREEN)

# Title
txb(s1, Inches(1.5), Inches(1.85), Inches(7.0), Inches(0.85),
    "2026 Paid Media Strategy",
    size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Subtitle
txb(s1, Inches(1.5), Inches(2.78), Inches(7.0), Inches(0.44),
    "Moving Budget to the Deals That Matter",
    size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# Meta
txb(s1, Inches(1.5), Inches(3.3), Inches(7.0), Inches(0.32),
    "Chris  ·  Head of Growth and Marketing  ·  April 2026",
    size=10, color=MUTED, align=PP_ALIGN.CENTER)


# ── SLIDE 2 — The Problem ─────────────────────────────────────────────────────
s2 = blank(prs)
bg(s2, WHITE)
top_bar(s2)
header(s2, "Strategic Context", "Notary Leads Cost More Than They Bring In")
footer(s2)

bullet_card(s2,
    l=Inches(0.42), t=Inches(1.0), w=Inches(4.52), h=Inches(4.0),
    header="Current State: The Problem",
    items=[
        "$5K/month goes to individual notary searches on Google",
        "Low value per customer. Most pay once and don't come back.",
        "Each deal takes a sales rep and Gina's full onboarding time. Cost likely exceeds revenue.",
        "No LinkedIn presence. B2B buyers who aren't searching are out of reach.",
        "RON leads are the whole pipeline. That won't scale.",
    ],
    bg_color=INK,
    text_color=LIGHT_BODY,
)

bullet_card(s2,
    l=Inches(5.06), t=Inches(1.0), w=Inches(4.52), h=Inches(4.0),
    header="Proposed State: The Shift",
    items=[
        "Same $5K/month. No budget increase needed.",
        "Reach FI ops, IT, and compliance buyers. Higher value and recurring revenue.",
        "Reach ISV developers through the Flex API. One integration can multiply your transaction base.",
        "Add LinkedIn to reach B2B buyers before they start searching.",
        "Keep RON only for institutional buyers.",
    ],
    bg_color=CANVAS,
    text_color=BODY,
)


# ── SLIDE 3 — Budget ──────────────────────────────────────────────────────────
s3 = blank(prs)
bg(s3, WHITE)
top_bar(s3)
header(s3, "Budget", "$5,000/Month, Moved to Enterprise Buyers")
footer(s3)

stat_data = [
    ("$3,000", "GOOGLE SEARCH ADS",
     "Five keyword groups: competitor displacement, industry use-case, "
     "authentication, Flex API, and institutional RON. Aimed at FI buyers "
     "and ISV developers."),
    ("$1,500", "LINKEDIN SPONSORED",
     "Two groups: FI ops and compliance leaders ($800) and ISV product teams "
     "($700). Sponsored posts only for the first two months."),
    ("$500",   "TEST & RESERVE",
     "Used for landing page tests, new keyword trials, and cost-per-click "
     "buffer. Moved to the top channel after a 30-day review."),
]
bw  = Inches(2.8)
gap = Inches(0.35)
bl  = Inches(0.65)

for i, (amount, label, desc) in enumerate(stat_data):
    x  = bl + i * (bw + gap)
    by = Inches(1.08)
    bh = Inches(2.85)
    rect(s3, x, by, bw, bh, CANVAS)
    rect(s3, x, by + bh - Inches(0.04), bw, Inches(0.04), GREEN)
    txb(s3, x, by + Inches(0.15), bw, Inches(0.6),
        amount, size=28, bold=True, color=INK, align=PP_ALIGN.CENTER)
    txb(s3, x, by + Inches(0.75), bw, Inches(0.25),
        label, size=7.5, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    txb(s3, x + Inches(0.12), by + Inches(1.05), bw - Inches(0.24), Inches(1.65),
        desc, size=10, color=BODY)

# Note box
ny = Inches(4.05)
rect(s3, Inches(0.65), ny, Inches(8.7), Inches(0.76), CANVAS)
rect(s3, Inches(0.65), ny, Inches(0.045), Inches(0.76), GREEN)
mixed_txb(s3, Inches(0.75), ny + Inches(0.12), Inches(8.45), Inches(0.55),
    segments=[
        ("No budget increase needed.  ", 10.5, True,  INK,  False),
        ("We are moving existing spend, not adding to it. "
         "Notary keywords stop right away. RON stays at a small budget "
         "for institutional buyers only.", 10.5, False, BODY, False),
    ]
)


# ── SLIDE 4 — Channels & Keywords ─────────────────────────────────────────────
s4 = blank(prs)
bg(s4, WHITE)
top_bar(s4)
header(s4, "Channels & Keywords", "Google + LinkedIn: Reaching the Right Buyers")
footer(s4)

# Keyword table
kw_rows = [
    ("Competitor Displacement", "~$900",  "docusign alternative for banks, adobe sign alternative compliance"),
    ("Industry / Use-Case",     "~$1,000","e-signature for credit unions, loan document electronic signature, e-sig NCUA compliant"),
    ("Authentication",          "~$600",  "e-signature identity verification, KBA electronic signature, e-sig fraud prevention"),
    ("Flex API / Developer",    "~$500",  "e-signature API financial services, white label e-signature API, embedded e-sign API"),
    ("RON Institutional",       "~$300",  "RON platform financial institutions — individual notary terms dropped"),
]
col_w = [Inches(1.45), Inches(0.55), Inches(2.35)]
tl    = Inches(0.42)
tt    = Inches(1.02)
tw    = sum(col_w)
th    = Inches(4.0)
row_h = Inches(0.52)

tbl_shape = s4.shapes.add_table(len(kw_rows) + 1, 3, tl, tt, tw, th)
tbl = tbl_shape.table

hdrs = ["Google Bucket", "Alloc.", "Top Keywords"]
for ci, h in enumerate(hdrs):
    cell = tbl.cell(0, ci)
    cell.fill.solid()
    cell.fill.fore_color.rgb = INK
    cell.text = h
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name  = FONT
    run.font.size  = Pt(8)
    run.font.bold  = True
    run.font.color.rgb = WHITE

for ri, (bucket, alloc, keys) in enumerate(kw_rows):
    row_bg = CANVAS if ri % 2 == 0 else WHITE
    for ci, val in enumerate([bucket, alloc, keys]):
        cell = tbl.cell(ri + 1, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = row_bg
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name  = FONT
        run.font.size  = Pt(8.5)
        run.font.bold  = (ci == 0)
        run.font.color.rgb = INK if ci == 0 else (GREEN if ci == 1 else BODY)

# Column widths
for ri in range(len(kw_rows) + 1):
    for ci, cw in enumerate(col_w):
        tbl.cell(ri, ci).width = cw

# LinkedIn cards
cx = Inches(5.0)
bullet_card(s4,
    l=cx, t=Inches(1.02), w=Inches(4.56), h=Inches(1.9),
    header="LinkedIn — FI Audience ($800)",
    items=[
        "VP Operations, IT Director, Director of Lending, CCO",
        "Banks and credit unions — 50–5,000 employees",
        "Show them how SIGNiX stops fraud and proves compliance",
    ],
    bg_color=CANVAS, text_color=BODY,
)
bullet_card(s4,
    l=cx, t=Inches(3.06), w=Inches(4.56), h=Inches(1.9),
    header="LinkedIn — ISV / Flex API ($700)",
    items=[
        "Product Manager, VP of Product, CTO, VP Engineering",
        "Fintech, LOS platforms, banking software — 10–500 employees",
        "Show them a white-label API built for regulated industries",
    ],
    bg_color=CANVAS, text_color=BODY,
)


# ── SLIDE 5 — Measurement ─────────────────────────────────────────────────────
s5 = blank(prs)
bg(s5, WHITE)
top_bar(s5)
header(s5, "Measurement", "60-Day Leading Signals  ·  90-Day Revenue Proof")
footer(s5)

leading = [
    ("SIGNAL", "Impression share above 40% on our top keywords"),
    ("SIGNAL", "CTR above 3–6% on B2B search campaigns"),
    ("SIGNAL", "Landing page bounce rate below 60%"),
    ("SIGNAL", "Form fill / demo request rate above 2% of paid clicks"),
    ("SIGNAL", "LinkedIn engagement rate above 0.5% sponsored content"),
    ("SIGNAL", "More than 40% of leads match our target buyer"),
    ("SIGNAL", "Lead to first sales touchpoint in under 48 hours"),
]
lagging = [
    ("PROOF", "Cost per qualified lead (CPL) vs. notary lead baseline"),
    ("PROOF", "MQLs per month (pipeline growth)"),
    ("PROOF", "Sales-accepted leads (SALs): did the team agree they were good?"),
    ("PROOF", "Pipeline sourced ($) attributed in HubSpot"),
    ("PROOF", "Average deal ACV vs. notary lead ACV"),
    ("PROOF", "Staff hours per closed deal: the number that shows real cost"),
    ("PROOF", "Transaction volume from new customers from paid ads"),
]

def indicator_col(slide, l, t, w, h, title, items, tag_color):
    # Section heading
    txb(slide, l, t, w, Inches(0.22), title.upper(),
        size=7.5, bold=True, color=GREEN)
    rect(slide, l, t + Inches(0.22), w, Pt(1.5), GREEN)

    iy = t + Inches(0.28)
    row_h = Inches(0.47)
    for tag, text in items:
        # tag pill
        rect(slide, l, iy + Inches(0.06), Inches(0.52), Inches(0.22), tag_color)
        txb(slide, l, iy + Inches(0.055), Inches(0.52), Inches(0.22),
            tag, size=6.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # text
        txb(slide, l + Inches(0.58), iy, w - Inches(0.6), Inches(0.42),
            text, size=9.5, color=BODY)
        rect(slide, l, iy + row_h - Pt(0.5), w, Pt(0.5), RULE)
        iy += row_h

indicator_col(s5, Inches(0.42), Inches(1.0), Inches(4.38), Inches(4.0),
              "Leading Indicators — Weekly", leading, GREEN)
indicator_col(s5, Inches(5.2), Inches(1.0), Inches(4.38), Inches(4.0),
              "Lagging Indicators — Monthly / Quarterly", lagging, INK)


# ── Save ──────────────────────────────────────────────────────────────────────
out = ("/Users/chris/Desktop/AI Summit 2026 Q2/"
       "proj-template-and-lease-SIGNiX-app/PROJECT-DOCS/"
       "SIGNIX-PAID-MEDIA-SLIDES.pptx")
prs.save(out)
print(f"Saved: {out}")
