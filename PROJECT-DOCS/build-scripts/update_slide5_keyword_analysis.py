#!/usr/bin/env python3
"""
Replace Slide 5 in SIGNiX_PaidMedia_4.6.26.pptx with the updated
keyword testing results slide (removed + validated + active buckets).
Leaves all other slides — including Chris's updated Slide 4 — untouched.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ── Colors ────────────────────────────────────────────────────────────────────
GREEN  = RGBColor(0x6d, 0xa3, 0x4a)
INK    = RGBColor(0x2e, 0x34, 0x40)
BODY   = RGBColor(0x54, 0x54, 0x54)
MUTED  = RGBColor(0x6b, 0x72, 0x80)
WHITE  = RGBColor(0xff, 0xff, 0xff)
CANVAS = RGBColor(0xf8, 0xfa, 0xfb)
RULE   = RGBColor(0xd8, 0xde, 0xe9)
RED    = RGBColor(0xc0, 0x39, 0x2b)
AMBER  = RGBColor(0xe6, 0x7e, 0x22)

FONT = "Calibri"
W    = Inches(10)
H    = Inches(5.625)

# ── Helpers ───────────────────────────────────────────────────────────────────

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def txb(slide, l, t, w, h, text, size=10, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    color = color or BODY
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name   = FONT
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def mixed_txb(slide, l, t, w, h, segments, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    for text, size, bold, color, italic in segments:
        run = p.add_run()
        run.text = text
        run.font.name   = FONT
        run.font.size   = Pt(size)
        run.font.bold   = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box

def top_bar(slide):
    rect(slide, Inches(0), Inches(0), W, Inches(0.055), GREEN)

def header(slide, eyebrow, title, title_size=18):
    txb(slide, Inches(0.65), Inches(0.09), Inches(8.7), Inches(0.22),
        eyebrow.upper(), size=7.5, bold=True, color=GREEN)
    txb(slide, Inches(0.65), Inches(0.31), Inches(8.7), Inches(0.55),
        title, size=title_size, bold=True, color=INK)

def footer(slide):
    fy = H - Inches(0.46)
    rect(slide, Inches(0), fy, W, Inches(0.46), WHITE)
    s = slide.shapes.add_shape(1, Inches(0), fy, W, Pt(0.75))
    s.fill.solid(); s.fill.fore_color.rgb = RULE; s.line.fill.background()
    txb(slide, W - Inches(1.55), fy + Inches(0.1),
        Inches(0.9), Inches(0.28),
        "SIGNiX", size=12, bold=True, color=INK, align=PP_ALIGN.RIGHT)
    rect(slide, W - Inches(0.6), fy + Inches(0.14),
         Inches(0.055), Inches(0.055), GREEN)

def section_label(slide, l, t, w, text):
    """Green all-caps section label with underline rule."""
    txb(slide, l, t, w, Inches(0.22), text.upper(), size=7.5, bold=True, color=GREEN)
    rect(slide, l, t + Inches(0.22), w, Pt(1.2), GREEN)

def anchor_box(slide, l, t, w, h, keyword, volume, trend, cpc, note):
    """Highlighted anchor term callout box."""
    rect(slide, l, t, w, h, CANVAS)
    rect(slide, l, t, Inches(0.045), h, GREEN)
    # Keyword
    txb(slide, l + Inches(0.1), t + Inches(0.08), w - Inches(0.15), Inches(0.22),
        keyword, size=9, bold=True, color=INK)
    # Stats row
    mixed_txb(slide, l + Inches(0.1), t + Inches(0.3), w - Inches(0.15), Inches(0.22),
        segments=[
            (volume + "  ", 9, True, GREEN, False),
            (trend + "   ", 8, False, BODY, False),
            ("CPC " + cpc, 8, False, MUTED, False),
        ])
    # Note
    txb(slide, l + Inches(0.1), t + Inches(0.52), w - Inches(0.15), Inches(0.26),
        note, size=8, color=MUTED, italic=True)


# ── Remove Slide 5 from existing file ────────────────────────────────────────
src = "/Users/chris/Desktop/SIGNiX_PaidMedia_4.6.26.pptx"
prs = Presentation(src)

slide_id_lst = prs.slides._sldIdLst
slide_id_el  = slide_id_lst[4]                         # 0-based index 4 = slide 5
r_id         = slide_id_el.get(qn("r:id"))
slide_id_lst.remove(slide_id_el)
# Remove the slide part from the package relationships
prs.part.drop_rel(r_id)

# ── Add new Slide 5 ───────────────────────────────────────────────────────────
new_slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(new_slide, WHITE)
top_bar(new_slide)
header(new_slide, "Keyword Testing Results — April 2026",
       "What We Kept, What We Cut, and Why")
footer(new_slide)

# ─────────────────────────────────────────────────────────────────────────────
# LEFT PANEL: Removed keywords table
# ─────────────────────────────────────────────────────────────────────────────
lx = Inches(0.42)
ct = Inches(1.02)

section_label(new_slide, lx, ct, Inches(4.5), "Removed Keywords — Zero Search Volume")

zero_rows = [
    ("docusign alternative for banks",             "0",  "Sector-specific term — no one searches this phrase"),
    ("adobe sign alternative compliance",           "0",  "Sector-specific term — zero volume confirmed"),
    ("sec 17a-4 e-signature",                       "0",  "Too specific; exam-triggered buyers use other paths"),
    ("finra digital signature",                     "0",  "Zero volume — compliance buyers don't search this"),
    ("electronic recordkeeping broker dealer",      "0",  "Zero volume — industry jargon, not buyer language"),
    ("ron platform compliance",                     "0",  "Zero volume — too niche as a search phrase"),
    ("remote online notary for banks",              "0",  "Zero volume — replaced with broader RON terms"),
    ("e signature for banks",                      "20", "Dying term — down 75% year over year"),
    ("digital signature mortgage docs",             "0",  "Zero volume + mortgage sector removed"),
    ("e signature rollover forms",                  "0",  "Buyers search document type, not signing action"),
    ("ira distribution e sign",                     "0",  "Industry shorthand — not how buyers search"),
    ("kba electronic signature",                    "0",  "KBA is internal jargon; buyers don't know the acronym"),
    ("e signature identity verification",           "0",  "Wrong word order — replaced with correct phrasing"),
    ("ai fraud e signature",                        "0",  "No search behavior yet. Re-test Q3 2026."),
    ("ron platform financial institutions",         "0",  "Zero volume — original RON institutional term"),
]

col_w = [Inches(1.55), Inches(0.32), Inches(2.48)]
tbl_t = ct + Inches(0.28)
tbl_h = Inches(3.6)

ztbl = new_slide.shapes.add_table(
    len(zero_rows) + 1, 3, lx, tbl_t, sum(col_w), tbl_h
).table

z_hdrs = ["Keyword", "Vol/Mo", "Reason Removed"]
for ci, h in enumerate(z_hdrs):
    cell = ztbl.cell(0, ci)
    cell.fill.solid(); cell.fill.fore_color.rgb = INK
    cell.text = h
    p = cell.text_frame.paragraphs[0]
    run = p.runs[0]
    run.font.name = FONT; run.font.size = Pt(7.5)
    run.font.bold = True; run.font.color.rgb = WHITE

for ri, (kw, vol, reason) in enumerate(zero_rows):
    bg_c = CANVAS if ri % 2 == 0 else WHITE
    for ci, val in enumerate([kw, vol, reason]):
        cell = ztbl.cell(ri + 1, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = bg_c
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        run = p.runs[0]
        run.font.name = FONT; run.font.size = Pt(7.5)
        # Volume column: red if 0 or small
        if ci == 1:
            run.font.color.rgb = RED if vol in ("0", "20") else BODY
            run.font.bold = True
        elif ci == 0:
            run.font.color.rgb = INK
        else:
            run.font.color.rgb = MUTED

for ri in range(len(zero_rows) + 1):
    for ci, cw in enumerate(col_w):
        ztbl.cell(ri, ci).width = cw

# Wrong buyer label below the table
wb_y = tbl_t + tbl_h + Inches(0.08)
rect(new_slide, lx, wb_y, Inches(4.35), Inches(0.28), CANVAS)
rect(new_slide, lx, wb_y, Inches(0.045), Inches(0.28), AMBER)
mixed_txb(new_slide, lx + Inches(0.1), wb_y + Inches(0.05),
          Inches(4.2), Inches(0.22),
    segments=[
        ("Also excluded — wrong buyer signal:  ", 7.5, True, INK, False),
        ("free docusign alternatives (1,900/mo), free/cheap variants (4 terms). "
         "Price-sensitive buyers. Not SIGNiX's customer.", 7.5, False, MUTED, False),
    ])

# ─────────────────────────────────────────────────────────────────────────────
# RIGHT PANEL: Validated anchors + active buckets
# ─────────────────────────────────────────────────────────────────────────────
rx = Inches(5.1)
ct2 = Inches(1.02)

section_label(new_slide, rx, ct2, Inches(4.5), "Validated Anchor Terms (Google Keyword Planner Confirmed)")

anchor_box(new_slide,
    l=rx, t=ct2 + Inches(0.3), w=Inches(4.5), h=Inches(0.88),
    keyword="beneficiary change form",
    volume="390/mo", trend="+50% YoY  ·  +23% last 3 mo",
    cpc="$1.02–$3.93",
    note="Best ROI keyword in the plan. Low competition. High identity-verification moment.")

anchor_box(new_slide,
    l=rx, t=ct2 + Inches(1.28), w=Inches(4.5), h=Inches(0.88),
    keyword="HIPAA compliant e-signature",
    volume="210/mo", trend="+191% last 3 mo  ·  +52% YoY",
    cpc="$11.07–$53.94",
    note="Fastest-growing term in the plan. Already in account. Anchor for healthcare bucket.")

anchor_box(new_slide,
    l=rx, t=ct2 + Inches(2.26), w=Inches(4.5), h=Inches(0.88),
    keyword="docusign competitors",
    volume="1,000/mo", trend="–32% YoY  ·  flat 3-mo",
    cpc="$13.00–$68.14",
    note="Best intent term in competitor displacement. Already in account. Medium competition.")

# Revised buckets summary
bkt_y = ct2 + Inches(3.24)
section_label(new_slide, rx, bkt_y, Inches(4.5), "Revised Budget Allocation — Based on Volume Data")

bkt_rows = [
    ("Compliance / Regulatory",  "~$400", "Reduced — zero volume confirmed; keep presence for rare high-intent searches"),
    ("Authentication / Fraud",   "~$500", "Reduced — testing new workflow-specific terms"),
    ("Wealth Management",        "~$700", "Increased — anchored by beneficiary change form"),
    ("Industry / Use-Case",      "~$800", "Increased — HIPAA growing fast; Legal/Debt Collection added"),
    ("Competitor Displacement",  "~$200", "Reduced — only 2 viable terms remain"),
    ("RON Institutional",        "~$100", "Hold — maintenance only"),
    ("Reserve / Test",           "~$300", "New — tests new Authentication and Legal/Debt terms"),
]

bkt_tbl = new_slide.shapes.add_table(
    len(bkt_rows) + 1, 3,
    rx, bkt_y + Inches(0.28),
    Inches(4.5), Inches(1.72)
).table

b_hdrs = ["Bucket", "$", "Rationale"]
for ci, h in enumerate(b_hdrs):
    cell = bkt_tbl.cell(0, ci)
    cell.fill.solid(); cell.fill.fore_color.rgb = INK
    cell.text = h
    p = cell.text_frame.paragraphs[0]
    run = p.runs[0]
    run.font.name = FONT; run.font.size = Pt(7.5)
    run.font.bold = True; run.font.color.rgb = WHITE

bkt_col_w = [Inches(1.45), Inches(0.42), Inches(2.63)]
for ri, (bucket, alloc, note) in enumerate(bkt_rows):
    bg_c = CANVAS if ri % 2 == 0 else WHITE
    for ci, val in enumerate([bucket, alloc, note]):
        cell = bkt_tbl.cell(ri + 1, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = bg_c
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        run = p.runs[0]
        run.font.name = FONT; run.font.size = Pt(7.5)
        if ci == 0:
            run.font.bold = True; run.font.color.rgb = INK
        elif ci == 1:
            run.font.bold = True; run.font.color.rgb = GREEN
        else:
            run.font.color.rgb = MUTED

for ri in range(len(bkt_rows) + 1):
    for ci, cw in enumerate(bkt_col_w):
        bkt_tbl.cell(ri, ci).width = cw

# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/chris/Desktop/SIGNiX_PaidMedia_4.6.26.pptx"
prs.save(out)
print(f"Saved: {out}")
