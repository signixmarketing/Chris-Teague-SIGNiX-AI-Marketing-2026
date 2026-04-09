#!/usr/bin/env python3
"""Build SIGNiX_PaidMedia_4.6.26.pptx — Rev 3, April 6 2026."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colors ────────────────────────────────────────────────────────────────────
GREEN  = RGBColor(0x6d, 0xa3, 0x4a)
INK    = RGBColor(0x2e, 0x34, 0x40)
BODY   = RGBColor(0x54, 0x54, 0x54)
MUTED  = RGBColor(0x6b, 0x72, 0x80)
WHITE  = RGBColor(0xff, 0xff, 0xff)
CANVAS = RGBColor(0xf8, 0xfa, 0xfb)
PANEL  = RGBColor(0xec, 0xef, 0xf4)
RULE   = RGBColor(0xd8, 0xde, 0xe9)
LIGHT_BODY = RGBColor(0xc8, 0xd0, 0xd8)
RED    = RGBColor(0xc0, 0x39, 0x2b)
AMBER  = RGBColor(0xe6, 0x7e, 0x22)

FONT = "Calibri"
W    = Inches(10)
H    = Inches(5.625)

# ── Core helpers ──────────────────────────────────────────────────────────────

def prs_new():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color, border_color=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if border_color:
        s.line.color.rgb = border_color
    else:
        s.line.fill.background()
    return s

def txb(slide, l, t, w, h, text, size=11, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    color = color or BODY
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = FONT
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def mixed_txb(slide, l, t, w, h, segments, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    for text, size, bold, color, italic in segments:
        run = p.add_run()
        run.text = text
        run.font.name  = FONT
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box

def bullet_card(slide, l, t, w, h, header, items,
                bg_color=CANVAS, text_color=BODY):
    rect(slide, l, t, w, h, bg_color)
    rect(slide, l, t, Inches(0.045), h, GREEN)

    box = slide.shapes.add_textbox(l + Inches(0.12), t + Inches(0.14),
                                   w - Inches(0.2), h - Inches(0.22))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(4)
    run = p.add_run()
    run.text = header.upper()
    run.font.name  = FONT
    run.font.size  = Pt(8.5)
    run.font.bold  = True
    run.font.color.rgb = GREEN

    for item in items:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = "›  " + item
        run.font.name  = FONT
        run.font.size  = Pt(10)
        run.font.color.rgb = text_color

def top_bar(slide):
    rect(slide, Inches(0), Inches(0), W, Inches(0.055), GREEN)

def header(slide, eyebrow, title, title_size=20):
    txb(slide, Inches(0.65), Inches(0.09), Inches(8.7), Inches(0.22),
        eyebrow.upper(), size=7.5, bold=True, color=GREEN)
    txb(slide, Inches(0.65), Inches(0.31), Inches(8.7), Inches(0.6),
        title, size=title_size, bold=True, color=INK)

def footer(slide):
    fy = H - Inches(0.46)
    rect(slide, Inches(0), fy, W, Inches(0.46), WHITE)
    rect(slide, Inches(0), fy, W, Pt(0.75), RULE)
    txb(slide, W - Inches(1.55), fy + Inches(0.1),
        Inches(0.9), Inches(0.28),
        "SIGNiX", size=12, bold=True, color=INK, align=PP_ALIGN.RIGHT)
    rect(slide, W - Inches(0.6), fy + Inches(0.14),
         Inches(0.055), Inches(0.055), GREEN)


# ── SLIDE 1 — Title ───────────────────────────────────────────────────────────
prs = prs_new()
s1  = blank(prs)
bg(s1, INK)

rect(s1, Inches(0), H - Inches(0.055), W, Inches(0.055), GREEN)

rect(s1, Inches(3.82), Inches(0.9), Inches(0.07), Inches(0.68), GREEN)
txb(s1, Inches(3.95), Inches(0.86), Inches(4.0), Inches(0.78),
    "SIGNiX", size=42, bold=True, color=WHITE)

rect(s1, Inches(3.82), Inches(1.66), Inches(4.0), Inches(0.022), GREEN)

txb(s1, Inches(1.5), Inches(1.85), Inches(7.0), Inches(0.85),
    "2026 Paid Media Strategy",
    size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txb(s1, Inches(1.5), Inches(2.78), Inches(7.0), Inches(0.44),
    "Updated April 6, 2026  ·  Rev 3",
    size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

txb(s1, Inches(1.5), Inches(3.3), Inches(7.0), Inches(0.32),
    "Chris  ·  Head of Growth and Marketing",
    size=10, color=MUTED, align=PP_ALIGN.CENTER)


# ── SLIDE 2 — The Problem ─────────────────────────────────────────────────────
s2 = blank(prs)
bg(s2, WHITE)
top_bar(s2)
header(s2, "Strategic Context", "Notary Leads Cost More Than They Bring In")
footer(s2)

bullet_card(s2,
    l=Inches(0.42), t=Inches(1.0), w=Inches(4.52), h=Inches(4.0),
    header="Current State: The Problem",
    items=[
        "Spend was concentrated on individual notary searches",
        "Low value per customer. High staff cost per deal.",
        "Each notary deal takes a sales rep and full CSM onboarding.",
        "Contribution margin per notary customer is near zero or negative.",
        "RON leads were the whole pipeline. That doesn't scale.",
    ],
    bg_color=INK,
    text_color=LIGHT_BODY,
)

bullet_card(s2,
    l=Inches(5.06), t=Inches(1.0), w=Inches(4.52), h=Inches(4.0),
    header="Proposed State: The Shift",
    items=[
        "Same $3,000/month. No budget increase.",
        "Three sectors: financial services / healthcare, wealth management, legal / debt collection.",
        "Every target sector has signing moments where ID Verify matters.",
        "Digital signature is legally stronger than a click — that's our story.",
        "Mortgage removed: too many barriers for new account acquisition.",
    ],
    bg_color=CANVAS,
    text_color=BODY,
)


# ── SLIDE 3 — Budget ──────────────────────────────────────────────────────────
s3 = blank(prs)
bg(s3, WHITE)
top_bar(s3)
header(s3, "Budget & Approach", "$3,000/Month — Google Only. LinkedIn Is Phase 2.")
footer(s3)

stat_data = [
    ("$3,000", "GOOGLE SEARCH ADS — ACTIVE",
     "Six keyword buckets: compliance/regulatory, authentication/fraud, "
     "wealth management, industry/use-case, competitor displacement, "
     "and RON institutional. Runs 8am–5pm EST Mon–Fri."),
    ("TBD", "LINKEDIN — PHASE 2",
     "On hold until 60-day Google review. Two planned audiences: "
     "FI ops/compliance leaders and ISV product teams. "
     "Launch target: June 2026."),
    ("60 days", "FIRST REVIEW",
     "CPL by bucket, CTR, lead qualification rate, and pipeline "
     "sourced in HubSpot. Lagging indicators (ACV, deals closed) "
     "carry a 90–180 day lag."),
]
bw  = Inches(2.8)
gap = Inches(0.35)
bl  = Inches(0.65)

for i, (amount, label, desc) in enumerate(stat_data):
    x  = bl + i * (bw + gap)
    by = Inches(1.08)
    bh = Inches(2.85)
    rect(s3, x, by, bw, bh, CANVAS)
    rect(s3, x, by + bh - Inches(0.04), bw, Inches(0.04), GREEN)
    txb(s3, x, by + Inches(0.15), bw, Inches(0.6),
        amount, size=28, bold=True, color=INK, align=PP_ALIGN.CENTER)
    txb(s3, x, by + Inches(0.75), bw, Inches(0.25),
        label, size=7.5, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    txb(s3, x + Inches(0.12), by + Inches(1.05), bw - Inches(0.24), Inches(1.65),
        desc, size=10, color=BODY)

ny = Inches(4.05)
rect(s3, Inches(0.65), ny, Inches(8.7), Inches(0.76), CANVAS)
rect(s3, Inches(0.65), ny, Inches(0.045), Inches(0.76), GREEN)
mixed_txb(s3, Inches(0.75), ny + Inches(0.12), Inches(8.45), Inches(0.55),
    segments=[
        ("No budget increase needed.  ", 10.5, True,  INK,  False),
        ("We are moving existing spend, not adding to it. "
         "Notary keywords stop right away. RON stays at $100/month "
         "for institutional buyers only.", 10.5, False, BODY, False),
    ]
)


# ── SLIDE 4 — Budget & Keywords ───────────────────────────────────────────────
s4 = blank(prs)
bg(s4, WHITE)
top_bar(s4)
header(s4, "Budget & Keywords", "$3,000/Month Google — Six Keyword Buckets")
footer(s4)

kw_rows = [
    ("Compliance / Regulatory", "~$400",  "SEC 17a-4 e-signature, FINRA digital signature, electronic recordkeeping broker dealer"),
    ("Authentication / Fraud",  "~$500",  "identity verification e-signature, verified electronic signature, e-sig fraud prevention banking"),
    ("Wealth Management",       "~$700",  "beneficiary change form, power of attorney e-signature, investment account opening signature"),
    ("Industry / Use-Case",     "~$800",  "HIPAA compliant e-signature, FDCPA compliant e-signature, wire transfer authorization form"),
    ("Competitor Displacement", "~$200",  "Docusign competitors, Docusign alternatives"),
    ("RON Institutional",       "~$100",  "RON platform compliance, remote online notary for banks — individual notary terms dropped"),
]
col_w = [Inches(1.5), Inches(0.55), Inches(2.3)]
tl    = Inches(0.42)
tt    = Inches(1.02)
tw    = sum(col_w)
th    = Inches(4.0)

tbl_shape = s4.shapes.add_table(len(kw_rows) + 1, 3, tl, tt, tw, th)
tbl = tbl_shape.table

hdrs = ["Google Bucket", "Alloc.", "Top Keywords"]
for ci, h in enumerate(hdrs):
    cell = tbl.cell(0, ci)
    cell.fill.solid()
    cell.fill.fore_color.rgb = INK
    cell.text = h
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name  = FONT
    run.font.size  = Pt(8)
    run.font.bold  = True
    run.font.color.rgb = WHITE

for ri, (bucket, alloc, keys) in enumerate(kw_rows):
    row_bg = CANVAS if ri % 2 == 0 else WHITE
    for ci, val in enumerate([bucket, alloc, keys]):
        cell = tbl.cell(ri + 1, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = row_bg
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name  = FONT
        run.font.size  = Pt(8.5)
        run.font.bold  = (ci == 0)
        run.font.color.rgb = INK if ci == 0 else (GREEN if ci == 1 else BODY)

for ri in range(len(kw_rows) + 1):
    for ci, cw in enumerate(col_w):
        tbl.cell(ri, ci).width = cw

# Right-side cards
cx = Inches(5.0)
bullet_card(s4,
    l=cx, t=Inches(1.02), w=Inches(4.56), h=Inches(1.9),
    header="Message Framework",
    items=[
        "Fear-first: 'Can you prove who signed?'",
        "Partnership close: 'We assign a CSM who knows your workflows.'",
        "Never lead with features: Lead with the consequence of not having it",
    ],
    bg_color=CANVAS, text_color=BODY,
)
bullet_card(s4,
    l=cx, t=Inches(3.06), w=Inches(4.56), h=Inches(1.9),
    header="Ad Scheduling + Phase 2",
    items=[
        "Ads run 8am–5pm EST, Mon–Fri only",
        "Leads arrive when someone can follow up same day",
        "LinkedIn added in Phase 2 after 60-day Google review",
    ],
    bg_color=CANVAS, text_color=BODY,
)


# ── SLIDE 5 — Keyword Testing Results (Removed + Validated + Active) ──────────
s5 = blank(prs)
bg(s5, WHITE)
top_bar(s5)
header(s5, "Keyword Testing Results — April 2026",
       "What We Kept, What We Cut, and Why", title_size=18)
footer(s5)

RED   = RGBColor(0xc0, 0x39, 0x2b)
AMBER = RGBColor(0xe6, 0x7e, 0x22)

def section_label(slide, l, t, w, text):
    txb(slide, l, t, w, Inches(0.22), text.upper(), size=7.5, bold=True, color=GREEN)
    rect(slide, l, t + Inches(0.22), w, Pt(1.2), GREEN)

def anchor_box(slide, l, t, w, h, keyword, volume, trend, cpc, note):
    rect(slide, l, t, w, h, CANVAS)
    rect(slide, l, t, Inches(0.045), h, GREEN)
    txb(slide, l + Inches(0.1), t + Inches(0.08), w - Inches(0.15), Inches(0.22),
        keyword, size=9, bold=True, color=INK)
    box = slide.shapes.add_textbox(l + Inches(0.1), t + Inches(0.3), w - Inches(0.15), Inches(0.22))
    box.word_wrap = True
    tf = box.text_frame
    p = tf.paragraphs[0]
    for seg_text, seg_size, seg_bold, seg_color, seg_italic in [
        (volume + "  ", 9, True, GREEN, False),
        (trend + "   ", 8, False, BODY, False),
        ("CPC " + cpc, 8, False, MUTED, False),
    ]:
        run = p.add_run()
        run.text = seg_text
        run.font.name = FONT; run.font.size = Pt(seg_size)
        run.font.bold = seg_bold; run.font.italic = seg_italic
        run.font.color.rgb = seg_color
    txb(slide, l + Inches(0.1), t + Inches(0.52), w - Inches(0.15), Inches(0.26),
        note, size=8, color=MUTED, italic=True)

# ── LEFT PANEL: Removed keywords ─────────────────────────────────────────────
lx  = Inches(0.42)
ct  = Inches(1.02)

section_label(s5, lx, ct, Inches(4.5), "Removed Keywords — Zero Search Volume")

zero_rows = [
    ("docusign alternative for banks",           "0",  "Sector-specific — no one searches this phrase"),
    ("adobe sign alternative compliance",         "0",  "Sector-specific — zero volume confirmed"),
    ("sec 17a-4 e-signature",                     "0",  "Too specific; exam-triggered buyers use other paths"),
    ("finra digital signature",                   "0",  "Compliance buyers don't search this term"),
    ("electronic recordkeeping broker dealer",    "0",  "Industry jargon — not how buyers search"),
    ("ron platform compliance",                   "0",  "Too niche as a search phrase"),
    ("remote online notary for banks",            "0",  "Zero volume — replaced with broader terms"),
    ("e signature for banks",                    "20", "Dying term — down 75% year over year"),
    ("digital signature mortgage docs",           "0",  "Zero volume + mortgage sector removed"),
    ("e signature rollover forms",                "0",  "Buyers search document type, not signing action"),
    ("ira distribution e sign",                   "0",  "Industry shorthand — not how buyers search"),
    ("kba electronic signature",                  "0",  "KBA is internal jargon; buyers don't know the acronym"),
    ("e signature identity verification",         "0",  "Wrong word order — replaced with correct phrasing"),
    ("ai fraud e signature",                      "0",  "No search behavior yet. Re-test Q3 2026."),
    ("ron platform financial institutions",       "0",  "Original RON term — zero volume confirmed"),
]

col_w = [Inches(1.55), Inches(0.32), Inches(2.48)]
tbl_t = ct + Inches(0.28)
tbl_h = Inches(3.6)

ztbl = s5.shapes.add_table(len(zero_rows) + 1, 3, lx, tbl_t, sum(col_w), tbl_h).table
z_hdrs = ["Keyword", "Vol/Mo", "Reason Removed"]
for ci, h in enumerate(z_hdrs):
    cell = ztbl.cell(0, ci)
    cell.fill.solid(); cell.fill.fore_color.rgb = INK
    cell.text = h
    p = cell.text_frame.paragraphs[0]; run = p.runs[0]
    run.font.name = FONT; run.font.size = Pt(7.5)
    run.font.bold = True; run.font.color.rgb = WHITE

for ri, (kw, vol, reason) in enumerate(zero_rows):
    bg_c = CANVAS if ri % 2 == 0 else WHITE
    for ci, val in enumerate([kw, vol, reason]):
        cell = ztbl.cell(ri + 1, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = bg_c
        cell.text = val
        p = cell.text_frame.paragraphs[0]; run = p.runs[0]
        run.font.name = FONT; run.font.size = Pt(7.5)
        if ci == 1:
            run.font.color.rgb = RED if vol in ("0", "20") else BODY
            run.font.bold = True
        elif ci == 0:
            run.font.color.rgb = INK
        else:
            run.font.color.rgb = MUTED

for ri in range(len(zero_rows) + 1):
    for ci, cw in enumerate(col_w):
        ztbl.cell(ri, ci).width = cw

# Wrong-buyer note
wb_y = tbl_t + tbl_h + Inches(0.08)
rect(s5, lx, wb_y, Inches(4.35), Inches(0.28), CANVAS)
s_wb = s5.shapes.add_shape(1, lx, wb_y, Inches(0.045), Inches(0.28))
s_wb.fill.solid(); s_wb.fill.fore_color.rgb = AMBER; s_wb.line.fill.background()
box_wb = s5.shapes.add_textbox(lx + Inches(0.1), wb_y + Inches(0.05), Inches(4.2), Inches(0.22))
box_wb.word_wrap = True
p_wb = box_wb.text_frame.paragraphs[0]
for seg_t, seg_s, seg_b, seg_c, seg_i in [
    ("Also excluded — wrong buyer signal:  ", 7.5, True, INK, False),
    ("free/cheap DocuSign variants (4 terms, 1,900+ searches). "
     "Price-sensitive buyers. Not SIGNiX's customer.", 7.5, False, MUTED, False),
]:
    r = p_wb.add_run(); r.text = seg_t
    r.font.name = FONT; r.font.size = Pt(seg_s)
    r.font.bold = seg_b; r.font.italic = seg_i
    r.font.color.rgb = seg_c

# ── RIGHT PANEL: Validated anchors + active buckets ──────────────────────────
rx  = Inches(5.1)
ct2 = Inches(1.02)

section_label(s5, rx, ct2, Inches(4.5), "Validated Anchor Terms (Google Keyword Planner Confirmed)")

anchor_box(s5, l=rx, t=ct2 + Inches(0.3), w=Inches(4.5), h=Inches(0.88),
    keyword="beneficiary change form",
    volume="390/mo", trend="+50% YoY  ·  +23% last 3 mo",
    cpc="$1.02–$3.93",
    note="Best ROI keyword in the plan. Low competition. High ID-Verify moment.")

anchor_box(s5, l=rx, t=ct2 + Inches(1.28), w=Inches(4.5), h=Inches(0.88),
    keyword="HIPAA compliant e-signature",
    volume="210/mo", trend="+191% last 3 mo  ·  +52% YoY",
    cpc="$11.07–$53.94",
    note="Fastest-growing term in the plan. Already in account. Healthcare anchor.")

anchor_box(s5, l=rx, t=ct2 + Inches(2.26), w=Inches(4.5), h=Inches(0.88),
    keyword="docusign competitors",
    volume="1,000/mo", trend="–32% YoY  ·  flat 3-mo",
    cpc="$13.00–$68.14",
    note="Best intent term in competitor displacement. Already in account.")

# Revised bucket summary
bkt_y = ct2 + Inches(3.24)
section_label(s5, rx, bkt_y, Inches(4.5), "Revised Allocation — Based on Volume Data")

bkt_rows = [
    ("Compliance / Regulatory",  "~$400", "Reduced — zero volume; keep for rare high-intent searches"),
    ("Authentication / Fraud",   "~$500", "Reduced — testing new workflow-specific terms"),
    ("Wealth Management",        "~$700", "Increased — anchored by beneficiary change form"),
    ("Industry / Use-Case",      "~$800", "Increased — HIPAA growing; Legal/Debt Collection added"),
    ("Competitor Displacement",  "~$200", "Reduced — only 2 viable terms remain"),
    ("RON Institutional",        "~$100", "Hold — maintenance only"),
]

bkt_tbl = s5.shapes.add_table(
    len(bkt_rows) + 1, 3, rx, bkt_y + Inches(0.28), Inches(4.5), Inches(1.44)
).table
b_hdrs = ["Bucket", "$", "Rationale"]
for ci, h in enumerate(b_hdrs):
    cell = bkt_tbl.cell(0, ci)
    cell.fill.solid(); cell.fill.fore_color.rgb = INK
    cell.text = h
    p = cell.text_frame.paragraphs[0]; run = p.runs[0]
    run.font.name = FONT; run.font.size = Pt(7.5)
    run.font.bold = True; run.font.color.rgb = WHITE

bkt_col_w = [Inches(1.45), Inches(0.42), Inches(2.63)]
for ri, (bucket, alloc, note) in enumerate(bkt_rows):
    bg_c = CANVAS if ri % 2 == 0 else WHITE
    for ci, val in enumerate([bucket, alloc, note]):
        cell = bkt_tbl.cell(ri + 1, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = bg_c
        cell.text = val
        p = cell.text_frame.paragraphs[0]; run = p.runs[0]
        run.font.name = FONT; run.font.size = Pt(7.5)
        if ci == 0:
            run.font.bold = True; run.font.color.rgb = INK
        elif ci == 1:
            run.font.bold = True; run.font.color.rgb = GREEN
        else:
            run.font.color.rgb = MUTED

for ri in range(len(bkt_rows) + 1):
    for ci, cw in enumerate(bkt_col_w):
        bkt_tbl.cell(ri, ci).width = cw


# ── SLIDE 6 — Measurement ─────────────────────────────────────────────────────
s6 = blank(prs)
bg(s6, WHITE)
top_bar(s6)
header(s6, "Measurement", "60-Day Leading Signals  ·  90-Day Revenue Proof")
footer(s6)

leading = [
    ("SIGNAL", "Impression share above 40% on top keyword buckets"),
    ("SIGNAL", "CTR above 3–6% on B2B search campaigns"),
    ("SIGNAL", "Landing page bounce rate below 60%"),
    ("SIGNAL", "Form fill / demo request rate above 2% of paid clicks"),
    ("SIGNAL", "More than 40% of leads match our target buyer"),
    ("SIGNAL", "Lead to first sales touchpoint in under 48 hours"),
]
lagging = [
    ("PROOF", "Cost per qualified lead (CPL) by bucket — which one works best?"),
    ("PROOF", "MQLs per month — is pipeline building?"),
    ("PROOF", "Sales-accepted leads (SALs) — did the team agree they were good?"),
    ("PROOF", "Pipeline sourced ($) attributed in HubSpot"),
    ("PROOF", "Average deal ACV vs. notary lead ACV baseline"),
    ("PROOF", "Transaction volume from new paid-sourced customers"),
]

def indicator_col(slide, l, t, w, h, title, items, tag_color):
    txb(slide, l, t, w, Inches(0.22), title.upper(),
        size=7.5, bold=True, color=GREEN)
    rect(slide, l, t + Inches(0.22), w, Pt(1.5), GREEN)

    iy = t + Inches(0.28)
    row_h = Inches(0.47)
    for tag, text in items:
        rect(slide, l, iy + Inches(0.06), Inches(0.52), Inches(0.22), tag_color)
        txb(slide, l, iy + Inches(0.055), Inches(0.52), Inches(0.22),
            tag, size=6.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txb(slide, l + Inches(0.58), iy, w - Inches(0.6), Inches(0.42),
            text, size=9.5, color=BODY)
        rect(slide, l, iy + row_h - Pt(0.5), w, Pt(0.5), RULE)
        iy += row_h

indicator_col(s6, Inches(0.42), Inches(1.0), Inches(4.38), Inches(4.0),
              "Leading Indicators — Weekly", leading, GREEN)
indicator_col(s6, Inches(5.2), Inches(1.0), Inches(4.38), Inches(4.0),
              "Lagging Indicators — Monthly / Quarterly", lagging, INK)


# ── Save ──────────────────────────────────────────────────────────────────────
out = ("/Users/chris/Desktop/AI Summit 2026 Q2/"
       "proj-template-and-lease-SIGNiX-app/PROJECT-DOCS/"
       "SIGNiX_PaidMedia_4.6.26.pptx")
prs.save(out)
print(f"Saved: {out}")
