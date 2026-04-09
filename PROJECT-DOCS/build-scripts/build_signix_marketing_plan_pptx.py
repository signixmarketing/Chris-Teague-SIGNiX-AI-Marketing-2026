#!/usr/bin/env python3
"""Build SIGNiX_MarketingPlan_4.7.26.pptx — 2026 Marketing Plan for leadership."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colors ─────────────────────────────────────────────────────────────────────
GREEN      = RGBColor(0x6d, 0xa3, 0x4a)
INK        = RGBColor(0x2e, 0x34, 0x40)
BODY       = RGBColor(0x54, 0x54, 0x54)
MUTED      = RGBColor(0x6b, 0x72, 0x80)
WHITE      = RGBColor(0xff, 0xff, 0xff)
CANVAS     = RGBColor(0xf8, 0xfa, 0xfb)
RULE       = RGBColor(0xd8, 0xde, 0xe9)
LIGHT_BODY = RGBColor(0xc8, 0xd0, 0xd8)
SLATE      = RGBColor(0x3a, 0x4a, 0x5a)   # funnel stage 2 — dark blue-gray
DARK_GREEN = RGBColor(0x4a, 0x72, 0x33)   # funnel stage 3 — deeper green

FONT = "Calibri"
W    = Inches(10)
H    = Inches(5.625)


# ── Core helpers ───────────────────────────────────────────────────────────────

def prs_new():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def txb(slide, l, t, w, h, text, size=11, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    color = color or BODY
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name   = FONT
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def mixed_txb(slide, l, t, w, h, segments, wrap=True):
    """Text box with multiple runs in one paragraph.
    segments = [(text, size, bold, color, italic)]"""
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    for text, size, bold, color, italic in segments:
        run = p.add_run()
        run.text = text
        run.font.name   = FONT
        run.font.size   = Pt(size)
        run.font.bold   = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box

def top_bar(slide):
    rect(slide, Inches(0), Inches(0), W, Inches(0.055), GREEN)

def slide_header(slide, eyebrow, title, title_size=20):
    txb(slide, Inches(0.65), Inches(0.09), Inches(8.7), Inches(0.22),
        eyebrow.upper(), size=7.5, bold=True, color=GREEN)
    txb(slide, Inches(0.65), Inches(0.31), Inches(8.7), Inches(0.6),
        title, size=title_size, bold=True, color=INK)

def footer(slide):
    fy = H - Inches(0.46)
    rect(slide, Inches(0), fy, W, Inches(0.46), WHITE)
    rect(slide, Inches(0), fy, W, Pt(0.75), RULE)
    txb(slide, W - Inches(1.55), fy + Inches(0.1),
        Inches(0.9), Inches(0.28),
        "SIGNiX", size=12, bold=True, color=INK, align=PP_ALIGN.RIGHT)
    rect(slide, W - Inches(0.6), fy + Inches(0.14),
         Inches(0.055), Inches(0.055), GREEN)


# ── Content helpers ────────────────────────────────────────────────────────────

def strategy_panel(slide, l, t, w, h, label, definition, bullets,
                   bg_color=CANVAS, label_color=INK, text_color=BODY):
    """Option B: Bold typographic anchor with definition line and bullets."""
    rect(slide, l, t, w, h, bg_color)
    rect(slide, l, t, Inches(0.045), h, GREEN)

    cx = l + Inches(0.18)
    cw = w - Inches(0.28)

    # Large typographic label
    txb(slide, cx, t + Inches(0.18), cw, Inches(0.70),
        label, size=44, bold=True, color=label_color)

    # Definition line
    txb(slide, cx, t + Inches(0.90), cw, Inches(0.26),
        definition, size=10.5, color=text_color)

    # Rule
    rect(slide, cx, t + Inches(1.20), cw, Pt(0.75), GREEN)

    # Bullet box
    box = slide.shapes.add_textbox(cx, t + Inches(1.30), cw, h - Inches(1.48))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = "\u203a  " + item
        run.font.name  = FONT
        run.font.size  = Pt(9.5)
        run.font.color.rgb = text_color


def belief_card(slide, l, t, w, h, segment, before_text, after_text):
    """Before/after belief card for the audience belief slide."""
    rect(slide, l, t, w, h, CANVAS)
    rect(slide, l, t, Inches(0.045), h, GREEN)

    box = slide.shapes.add_textbox(
        l + Inches(0.16), t + Inches(0.14),
        w - Inches(0.26), h - Inches(0.22)
    )
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True

    # Segment label
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = segment.upper()
    r0.font.name  = FONT
    r0.font.size  = Pt(8.5)
    r0.font.bold  = True
    r0.font.color.rgb = GREEN

    # TODAY line
    p1 = tf.add_paragraph()
    p1.space_before = Pt(6)
    r1a = p1.add_run()
    r1a.text = "TODAY:  "
    r1a.font.name  = FONT
    r1a.font.size  = Pt(7.5)
    r1a.font.bold  = True
    r1a.font.color.rgb = MUTED
    r1b = p1.add_run()
    r1b.text = before_text
    r1b.font.name   = FONT
    r1b.font.size   = Pt(9)
    r1b.font.italic = True
    r1b.font.color.rgb = MUTED

    # AFTER line
    p2 = tf.add_paragraph()
    p2.space_before = Pt(9)
    r2a = p2.add_run()
    r2a.text = "AFTER SIGNIX:  "
    r2a.font.name  = FONT
    r2a.font.size  = Pt(7.5)
    r2a.font.bold  = True
    r2a.font.color.rgb = GREEN
    r2b = p2.add_run()
    r2b.text = after_text
    r2b.font.name  = FONT
    r2b.font.size  = Pt(9.5)
    r2b.font.bold  = True
    r2b.font.color.rgb = INK


def audience_card(slide, l, t, w, h, icon, segment, who, behavior, channel):
    """Target audience card: icon anchor, segment label, who-they-are, behavior target, channel."""
    rect(slide, l, t, w, h, CANVAS)
    rect(slide, l, t, w, Inches(0.055), GREEN)   # green top stripe

    cx = l + Inches(0.20)
    cw = w - Inches(0.32)

    # Large icon
    txb(slide, cx, t + Inches(0.10), Inches(0.60), Inches(0.60),
        icon, size=30, color=GREEN)

    # Segment name
    txb(slide, cx + Inches(0.66), t + Inches(0.16), cw - Inches(0.66), Inches(0.36),
        segment.upper(), size=10, bold=True, color=INK)

    # Rule
    rect(slide, cx, t + Inches(0.76), cw, Pt(0.75), RULE)

    # Who they are
    txb(slide, cx, t + Inches(0.86), cw, Inches(0.28),
        who, size=8.5, italic=True, color=MUTED)

    # Behavior target label + text
    txb(slide, cx, t + Inches(1.18), cw, Inches(0.20),
        "WHAT WE WANT THEM TO DO", size=7, bold=True, color=GREEN)

    box = slide.shapes.add_textbox(cx, t + Inches(1.40), cw, h - Inches(1.58))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = behavior
    run.font.name  = FONT
    run.font.size  = Pt(9.5)
    run.font.bold  = True
    run.font.color.rgb = INK

    # Channel / reach note
    txb(slide, cx, t + h - Inches(0.28), cw, Inches(0.22),
        channel, size=7.5, color=MUTED)


def initiative_card(slide, l, t, w, h, number, name, growth_play, detail,
                    status_text, stripe_color=GREEN, status_color=GREEN):
    """Card grid item for the initiatives slide.
    stripe_color = visual status indicator at top of card."""
    rect(slide, l, t, w, h, CANVAS)

    # Colored top stripe
    rect(slide, l, t, w, Inches(0.10), stripe_color)

    # Large faded number — background visual anchor, top-right
    txb(slide, l + w - Inches(0.78), t + Inches(0.10), Inches(0.72), Inches(0.55),
        number, size=34, bold=True, color=RULE, align=PP_ALIGN.RIGHT)

    # Initiative name
    txb(slide, l + Inches(0.14), t + Inches(0.18), w - Inches(0.90), Inches(0.30),
        name, size=11, bold=True, color=INK)

    # Growth play label
    txb(slide, l + Inches(0.14), t + Inches(0.50), w - Inches(0.22), Inches(0.20),
        growth_play, size=8.5, color=MUTED)

    # Rule
    rect(slide, l + Inches(0.14), t + Inches(0.73), w - Inches(0.28), Pt(0.5), RULE)

    # Detail line
    txb(slide, l + Inches(0.14), t + Inches(0.77), w - Inches(0.22), Inches(0.20),
        detail, size=8.5, color=BODY)

    # Status dot + text pinned to bottom of card
    rect(slide, l + Inches(0.14), t + h - Inches(0.30),
         Inches(0.055), Inches(0.055), status_color)
    txb(slide, l + Inches(0.24), t + h - Inches(0.34), w - Inches(0.36), Inches(0.26),
        status_text, size=8, color=status_color)


def funnel_stage(slide, l, t, w, h, stage_label, big_number, big_sub,
                 initiatives, bg_color=INK, rule_color=GREEN):
    """Horizontal funnel stage: label, large metric, supporting initiatives."""
    rect(slide, l, t, w, h, bg_color)

    cx = l + Inches(0.18)
    cw = w - Inches(0.28)

    # Stage label
    txb(slide, cx, t + Inches(0.16), cw, Inches(0.20),
        stage_label.upper(), size=7.5, bold=True, color=LIGHT_BODY)

    # Big metric anchor
    txb(slide, cx, t + Inches(0.38), cw, Inches(0.60),
        big_number, size=34, bold=True, color=WHITE)

    # Metric description
    txb(slide, cx, t + Inches(0.98), cw, Inches(0.22),
        big_sub, size=9, italic=True, color=LIGHT_BODY)

    # Rule
    rect(slide, cx, t + Inches(1.22), cw, Pt(0.75), rule_color)

    # Initiative bullets
    box = slide.shapes.add_textbox(cx, t + Inches(1.32), cw, h - Inches(1.50))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(initiatives):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(7)
        run = p.add_run()
        run.text = "\u203a  " + item
        run.font.name  = FONT
        run.font.size  = Pt(8.5)
        run.font.color.rgb = LIGHT_BODY


def oval(slide, l, t, w, h, fill_color, stroke_color=None, stroke_width=Pt(1.5)):
    """Ellipse / circle shape."""
    s = slide.shapes.add_shape(9, l, t, w, h)   # 9 = OVAL
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if stroke_color:
        s.line.color.rgb = stroke_color
        s.line.width     = stroke_width
    else:
        s.line.fill.background()
    return s


def rnd_rect(slide, l, t, w, h, fill_color, stroke_color=None, stroke_width=Pt(1.5)):
    """Rounded rectangle shape."""
    s = slide.shapes.add_shape(5, l, t, w, h)   # 5 = ROUNDED_RECTANGLE
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if stroke_color:
        s.line.color.rgb = stroke_color
        s.line.width     = stroke_width
    else:
        s.line.fill.background()
    return s


def spoke(slide, x1, y1, x2, y2, color=RULE, width=Pt(1.5)):
    """Straight connector line between two points."""
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)   # 1 = STRAIGHT
    conn.line.color.rgb = color
    conn.line.width     = width
    return conn


def hub_node(slide, cx, cy, icon, seg_name, behavior):
    """Rounded card centered at (cx, cy) for the hub-and-spoke audience slide."""
    NW = Inches(2.30)
    NH = Inches(1.22)
    nl = cx - NW / 2
    nt = cy - NH / 2

    rnd_rect(slide, nl, nt, NW, NH, CANVAS, stroke_color=GREEN, stroke_width=Pt(1.5))
    rect(slide, nl, nt, NW, Inches(0.06), GREEN)

    txb(slide, nl + Inches(0.12), nt + Inches(0.08), Inches(0.40), Inches(0.40),
        icon, size=22, color=GREEN)
    txb(slide, nl + Inches(0.54), nt + Inches(0.12), NW - Inches(0.62), Inches(0.30),
        seg_name, size=9, bold=True, color=INK)
    rect(slide, nl + Inches(0.12), nt + Inches(0.52), NW - Inches(0.22), Pt(0.5), RULE)

    box = slide.shapes.add_textbox(
        nl + Inches(0.12), nt + Inches(0.58),
        NW - Inches(0.22), NH - Inches(0.68))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text           = behavior
    run.font.name      = FONT
    run.font.size      = Pt(8.5)
    run.font.color.rgb = BODY


def make_table(slide, rows, col_widths, tl, tt, th, headers,
               hdr_bg=INK, hdr_fg=WHITE):
    """Build a styled data table. rows = list of tuples."""
    tw = sum(col_widths)
    ncols = len(col_widths)
    shape = slide.shapes.add_table(len(rows) + 1, ncols, tl, tt, tw, th)
    tbl = shape.table

    # Header row
    for ci, h_text in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = hdr_bg
        cell.text = h_text
        p = cell.text_frame.paragraphs[0]
        run = p.runs[0]
        run.font.name  = FONT
        run.font.size  = Pt(8)
        run.font.bold  = True
        run.font.color.rgb = hdr_fg

    # Data rows
    for ri, row_data in enumerate(rows):
        row_bg = CANVAS if ri % 2 == 0 else WHITE
        for ci, val in enumerate(row_data):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_bg
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0]
            run.font.name  = FONT
            run.font.size  = Pt(8.5)
            run.font.bold  = (ci == 0)
            run.font.color.rgb = INK if ci == 0 else BODY

    # Column widths
    for ri in range(len(rows) + 1):
        for ci, cw in enumerate(col_widths):
            tbl.cell(ri, ci).width = cw


# ── SLIDE 1 — Title ────────────────────────────────────────────────────────────
prs = prs_new()
s1  = blank(prs)
bg(s1, INK)

rect(s1, Inches(0), H - Inches(0.055), W, Inches(0.055), GREEN)

rect(s1, Inches(3.82), Inches(0.9), Inches(0.07), Inches(0.68), GREEN)
txb(s1, Inches(3.95), Inches(0.86), Inches(4.0), Inches(0.78),
    "SIGNiX", size=42, bold=True, color=WHITE)

rect(s1, Inches(3.82), Inches(1.66), Inches(4.0), Inches(0.022), GREEN)

txb(s1, Inches(1.5), Inches(1.88), Inches(7.0), Inches(0.85),
    "2026 Marketing Plan",
    size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txb(s1, Inches(1.5), Inches(2.80), Inches(7.0), Inches(0.46),
    "A focused plan to grow transaction volume and win new accounts",
    size=13, color=GREEN, align=PP_ALIGN.CENTER)

txb(s1, Inches(1.5), Inches(3.36), Inches(7.0), Inches(0.32),
    "Chris Teague  |  Head of Growth and Marketing  |  April 2026",
    size=10, color=MUTED, align=PP_ALIGN.CENTER)


# ── SLIDE 2 — Two Strategies ───────────────────────────────────────────────────
s2 = blank(prs)
bg(s2, WHITE)
top_bar(s2)
slide_header(s2, "Strategic Framework", "Two Strategies. One Growth Plan.")
footer(s2)

strategy_panel(s2,
    l=Inches(0.42), t=Inches(1.02), w=Inches(4.38), h=Inches(3.90),
    label="DEPTH",
    definition="Get more from the customers we already have",
    bullets=[
        "Many TruStage FIs have SIGNiX available but use zero authentications. "
        "That is revenue we have not captured yet.",
        "Mike's fraud email and the TruStage CSM brief turn awareness "
        "into discovery calls.",
        "Goal: turn low-activity FI accounts into regular authentication users.",
    ],
    bg_color=INK,
    label_color=GREEN,
    text_color=LIGHT_BODY,
)

strategy_panel(s2,
    l=Inches(5.20), t=Inches(1.02), w=Inches(4.38), h=Inches(3.90),
    label="WIDTH",
    definition="Win new accounts in segments where we have a clear edge",
    bullets=[
        "We target five segments where a regulator or court requires signatures "
        "that hold up as verified proof.",
        "Google Ads, ABM, and integration partners like Gannett and Convoke "
        "are our three growth channels.",
        "Goal: new transaction volume from new accounts and partners by Q3 2026.",
    ],
    bg_color=CANVAS,
    label_color=INK,
    text_color=BODY,
)


# ── SLIDE 3 — Five Target Segments ────────────────────────────────────────────
s3 = blank(prs)
bg(s3, WHITE)
top_bar(s3)
slide_header(s3, "Target Audiences", "We Win Where Signatures Have to Stand Up")
footer(s3)

seg_rows = [
    ("Wealth Management",
     "Pershing, Osaic, Docupace",
     "Documents must hold up under FINRA and SEC review. "
     "A disputed wire or transfer that standard e-sign cannot defend."),
    ("Credit Unions and Banks",
     "TruStage FI network",
     "ID Verify proves who signed every loan document. "
     "A fraud case where a member denies signing triggers the switch."),
    ("Debt Buyers and Servicers",
     "Cavalry, InvestiNet, Convoke",
     "Chain of title and payment agreements must hold up in court. "
     "A challenged collection agreement drives the conversation."),
    ("Media and Legal Publishers",
     "Advance Local, Gannett",
     "Legal notices require notarization by state law. "
     "DocuSign and Adobe cannot do Remote Online Notary. SIGNiX can."),
    ("Software Partners",
     "Docupace, Convoke, Lineup",
     "Clients are asking for stronger identity proof than click-to-sign. "
     "SIGNiX gives partners a compliance-grade option competitors cannot match."),
]

make_table(
    s3, seg_rows,
    col_widths=[Inches(1.75), Inches(1.70), Inches(4.75)],
    tl=Inches(0.42), tt=Inches(1.02), th=Inches(3.95),
    headers=["Segment", "Examples", "What Drives the Move to SIGNiX"],
)


# ── SLIDE 4 — Hub-and-Spoke Audiences ─────────────────────────────────────────
s4 = blank(prs)
bg(s4, WHITE)
top_bar(s4)
slide_header(s4, "Who We Serve", "Four Audiences. One Growth Strategy.")
footer(s4)

HUB_X = Inches(5.00)
HUB_Y = Inches(3.06)

NODE_CENTERS = [
    (Inches(1.65), Inches(1.86)),   # top-left
    (Inches(8.35), Inches(1.86)),   # top-right
    (Inches(1.65), Inches(4.26)),   # bottom-left
    (Inches(8.35), Inches(4.26)),   # bottom-right
]

# Spokes drawn first so hub and node cards render on top
for nx, ny in NODE_CENTERS:
    spoke(s4, HUB_X, HUB_Y, nx, ny, color=RULE, width=Pt(1.5))

# Outer decorative ring
oval(s4, Inches(4.18), Inches(2.38), Inches(1.64), Inches(1.36),
     CANVAS, stroke_color=RULE, stroke_width=Pt(1.0))

# Inner hub circle
oval(s4, Inches(4.32), Inches(2.48), Inches(1.36), Inches(1.16), INK)

# Hub text
txb(s4, Inches(4.32), Inches(2.57), Inches(1.36), Inches(0.38),
    "SIGNiX", size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s4, Inches(4.32), Inches(2.96), Inches(1.36), Inches(0.22),
    "Flex API and MyDox", size=7, color=LIGHT_BODY, align=PP_ALIGN.CENTER)

# Segment node cards
hub_segments = [
    ("\u2696", "Wealth Management",
     "Turn on ID verification for every client document."),
    ("\u2302", "Credit Unions and Banks",
     "Add authentication to every loan signing event."),
    ("\u2261", "Debt, Legal and Media",
     "Use RON for any legal notice that needs notarizing."),
    ("\u25c6", "Integration Partners",
     "Embed SIGNiX as the compliance-grade option in your platform."),
]

for (nx, ny), (icon, seg, behavior) in zip(NODE_CENTERS, hub_segments):
    hub_node(s4, nx, ny, icon, seg, behavior)


# ── SLIDE 5 — Active Initiatives ──────────────────────────────────────────────
s5 = blank(prs)
bg(s5, WHITE)
top_bar(s5)
slide_header(s5, "Active Initiatives", "Six Things We Are Running Right Now")
footer(s5)

# Legend: GREEN stripe = Active  |  INK stripe = In progress  |  RULE stripe = Pending
cards = [
    # (number, name, growth_play, detail, status_text, stripe_color, status_color)
    ("01", "Google Ads",
     "Width: attract new accounts",
     "19 validated keywords across 3 sectors",
     "Active",
     GREEN, GREEN),
    ("02", "ABM Scorecard",
     "Width: earn share in key segments",
     "35 named accounts, rollout ready",
     "Built, rollout pending",
     RULE, MUTED),
    ("03", "Authentication Activation",
     "Depth: more revenue from existing FIs",
     "TruStage FI network",
     "CTO data request pending",
     RULE, MUTED),
    ("04", "Mike Fraud Email",
     "Depth: keep SIGNiX visible",
     "Monthly send, March complete",
     "Active",
     GREEN, GREEN),
    ("05", "Lineup / Gannett",
     "Width: media publishers and RON",
     "Replaces SignRequest in Lineup",
     "Integration in discussion",
     INK, INK),
    ("06", "Info-Tech Briefing",
     "Width: IT buyer visibility",
     "Analyst coverage secured at no cost",
     "Scheduling pending",
     RULE, MUTED),
]

# 3-column x 2-row grid
CARD_W = Inches(2.93)
CARD_H = Inches(1.95)
cols   = [Inches(0.42), Inches(3.535), Inches(6.65)]
rows   = [Inches(1.02), Inches(3.11)]

for idx, (num, name, play, detail, status, sc, stc) in enumerate(cards):
    col = cols[idx % 3]
    row = rows[idx // 3]
    initiative_card(s5, col, row, CARD_W, CARD_H,
                    num, name, play, detail, status, sc, stc)

# Legend strip across the bottom
legend_y = Inches(5.07)
for color, label in [(GREEN, "Active"), (INK, "In progress"), (RULE, "Pending")]:
    pass  # legend handled by stripe color — no additional strip needed


# ── SLIDE 6 — Funnel Measurement Dashboard ────────────────────────────────────
s6 = blank(prs)
bg(s6, WHITE)
top_bar(s6)
slide_header(s6, "Measurement", "What Success Looks Like at 60 Days")
footer(s6)

# Three funnel stages: AWARENESS → ENGAGEMENT → CONVERSION
# Heights decrease left to right; all vertically centered at y=3.06"
Y_CENTER = Inches(3.06)
COL_W    = Inches(3.053)

stages = [
    # (label, big_number, big_sub, [initiatives], bg_color, rule_color, height)
    ("Awareness",
     "15+",
     "Qualified leads per month",
     [
         "Google Ads — 19 keywords across 3 sectors",
         "Info-Tech — 1 analyst report published",
     ],
     INK, GREEN, Inches(3.80)),

    ("Engagement",
     "10",
     "Active accounts per week",
     [
         "ABM Scorecard — 1 or more touchpoints per week",
         "Mike Fraud Email — 5% open rate, 2 calls per send",
     ],
     SLATE, GREEN, Inches(3.10)),

    ("Conversion",
     "2",
     "Agreements or pilots closed",
     [
         "Authentication Activation — 1 FI pilot turned on",
         "Lineup / Gannett — integration agreement signed",
     ],
     DARK_GREEN, LIGHT_BODY, Inches(2.50)),
]

for i, (label, num, sub, inits, bg_col, rule_col, stage_h) in enumerate(stages):
    col_l = Inches(0.42) + i * COL_W
    col_t = Y_CENTER - stage_h / 2
    funnel_stage(s6, col_l, col_t, COL_W, stage_h,
                 label, num, sub, inits, bg_col, rule_col)

# Flow connectors between stages
for i in range(2):
    cx = Inches(0.42) + (i + 1) * COL_W - Inches(0.14)
    txb(s6, cx, Y_CENTER - Inches(0.14), Inches(0.28), Inches(0.28),
        "\u203a", size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Legend strip below stages
legend_t = Y_CENTER + Inches(1.95)
legend_items = [
    (INK,       "Awareness — top of funnel"),
    (SLATE,     "Engagement — mid funnel"),
    (DARK_GREEN,"Conversion — bottom of funnel"),
]
lx = Inches(0.42)
for lc, ltxt in legend_items:
    rect(s6, lx, legend_t, Inches(0.12), Inches(0.12), lc)
    txb(s6, lx + Inches(0.18), legend_t - Inches(0.02), Inches(2.50), Inches(0.18),
        ltxt, size=7.5, color=MUTED)
    lx += Inches(3.053)


# ── Save ───────────────────────────────────────────────────────────────────────
OUT = (
    "/Users/chris/Desktop/AI Summit 2026 Q2/"
    "proj-template-and-lease-SIGNiX-app/PROJECT-DOCS/"
    "SIGNiX_MarketingPlan_4.7.26.pptx"
)
prs.save(OUT)
print(f"Saved: {OUT}")
