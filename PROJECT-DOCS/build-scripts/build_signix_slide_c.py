#!/usr/bin/env python3
"""Build SIGNiX_SlideC_Audiences.pptx — Single slide, Option C hub-and-spoke layout."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colors ──────────────────────────────────────────────────────────────────
GREEN      = RGBColor(0x6d, 0xa3, 0x4a)
INK        = RGBColor(0x2e, 0x34, 0x40)
BODY       = RGBColor(0x54, 0x54, 0x54)
MUTED      = RGBColor(0x6b, 0x72, 0x80)
WHITE      = RGBColor(0xff, 0xff, 0xff)
CANVAS     = RGBColor(0xf8, 0xfa, 0xfb)
RULE       = RGBColor(0xd8, 0xde, 0xe9)
LIGHT_BODY = RGBColor(0xc8, 0xd0, 0xd8)

FONT = "Calibri"
W    = Inches(10)
H    = Inches(5.625)

# MSO shape type integers
SHAPE_RECT          = 1
SHAPE_OVAL          = 9
SHAPE_ROUNDED_RECT  = 5
CONNECTOR_STRAIGHT  = 1   # MSO_CONNECTOR_TYPE.STRAIGHT


# ── Helpers ─────────────────────────────────────────────────────────────────

def blank_slide(prs):
    prs.slide_width  = W
    prs.slide_height = H
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(SHAPE_RECT, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def oval(slide, l, t, w, h, fill_color, stroke_color=None, stroke_width=Pt(1.5)):
    s = slide.shapes.add_shape(SHAPE_OVAL, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if stroke_color:
        s.line.color.rgb   = stroke_color
        s.line.width       = stroke_width
    else:
        s.line.fill.background()
    return s


def rnd_rect(slide, l, t, w, h, fill_color, stroke_color=None, stroke_width=Pt(1.5)):
    s = slide.shapes.add_shape(SHAPE_ROUNDED_RECT, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if stroke_color:
        s.line.color.rgb   = stroke_color
        s.line.width       = stroke_width
    else:
        s.line.fill.background()
    return s


def txb(slide, l, t, w, h, text, size=11, bold=False, italic=False,
        color=None, align=PP_ALIGN.LEFT):
    color = color or BODY
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text            = text
    run.font.name       = FONT
    run.font.size       = Pt(size)
    run.font.bold       = bold
    run.font.italic     = italic
    run.font.color.rgb  = color
    return box


def spoke(slide, x1, y1, x2, y2, color=RULE, width=Pt(1.5)):
    """Straight connector line from hub center to node center."""
    conn = slide.shapes.add_connector(CONNECTOR_STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width     = width
    return conn


def top_bar(slide):
    rect(slide, Inches(0), Inches(0), W, Inches(0.055), GREEN)


def footer(slide):
    fy = H - Inches(0.46)
    rect(slide, Inches(0), fy, W, Inches(0.46), WHITE)
    rect(slide, Inches(0), fy, W, Pt(0.75), RULE)
    txb(slide, Inches(0.42), fy + Inches(0.10), Inches(4), Inches(0.28),
        "SIGNiX Confidential", size=6.5, color=MUTED)
    txb(slide, W - Inches(1.60), fy + Inches(0.10), Inches(1.20), Inches(0.28),
        "signix.com", size=6.5, color=MUTED, align=PP_ALIGN.RIGHT)


# ── Segment node helper ──────────────────────────────────────────────────────

def node_card(slide, cx, cy, icon, seg_name, behavior):
    """Draw a rounded card centered at (cx, cy) with icon, segment, and behavior."""
    NW = Inches(2.30)
    NH = Inches(1.22)
    nl = cx - NW / 2
    nt = cy - NH / 2

    # Card background + green stroke
    rnd_rect(slide, nl, nt, NW, NH, CANVAS, stroke_color=GREEN, stroke_width=Pt(1.5))

    # Green top stripe
    rect(slide, nl, nt, NW, Inches(0.06), GREEN)

    # Icon
    txb(slide, nl + Inches(0.12), nt + Inches(0.08), Inches(0.40), Inches(0.40),
        icon, size=22, color=GREEN)

    # Segment name
    txb(slide, nl + Inches(0.54), nt + Inches(0.12), NW - Inches(0.62), Inches(0.30),
        seg_name, size=9, bold=True, color=INK)

    # Thin rule
    rect(slide, nl + Inches(0.12), nt + Inches(0.52), NW - Inches(0.22), Pt(0.5), RULE)

    # Behavior text
    box = slide.shapes.add_textbox(
        nl + Inches(0.12), nt + Inches(0.58),
        NW - Inches(0.22), NH - Inches(0.68))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text           = behavior
    run.font.name      = FONT
    run.font.size      = Pt(8.5)
    run.font.color.rgb = BODY


# ── Build slide ──────────────────────────────────────────────────────────────

prs = Presentation()
s = blank_slide(prs)
bg(s, WHITE)
top_bar(s)
footer(s)

# Slide label and title
txb(s, Inches(0.65), Inches(0.09), Inches(8), Inches(0.22),
    "WHO WE SERVE", size=7.5, bold=True, color=GREEN)
txb(s, Inches(0.65), Inches(0.31), Inches(8), Inches(0.50),
    "Four Audiences. One Growth Strategy.", size=18, bold=True, color=INK)

# Hub center point
HUB_X = Inches(5.00)
HUB_Y = Inches(3.06)

# Node centers (top-left, top-right, bottom-left, bottom-right)
NODE_CENTERS = [
    (Inches(1.65), Inches(1.86)),
    (Inches(8.35), Inches(1.86)),
    (Inches(1.65), Inches(4.26)),
    (Inches(8.35), Inches(4.26)),
]

# ── Step 1: Spokes (draw before shapes so circles sit on top) ────────────────
for nx, ny in NODE_CENTERS:
    spoke(s, HUB_X, HUB_Y, nx, ny, color=RULE, width=Pt(1.5))

# ── Step 2: Center hub ───────────────────────────────────────────────────────
# Outer decorative ring
oval(s, Inches(4.18), Inches(2.38), Inches(1.64), Inches(1.36),
     CANVAS, stroke_color=RULE, stroke_width=Pt(1.0))

# Inner filled circle
oval(s, Inches(4.32), Inches(2.48), Inches(1.36), Inches(1.16), INK)

# Hub label
txb(s, Inches(4.32), Inches(2.57), Inches(1.36), Inches(0.38),
    "SIGNiX", size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, Inches(4.32), Inches(2.96), Inches(1.36), Inches(0.22),
    "Flex API", size=8, color=LIGHT_BODY, align=PP_ALIGN.CENTER)

# ── Step 3: Segment node cards ───────────────────────────────────────────────
segments = [
    ("\u2696", "Wealth Management",
     "Turn on ID verification for every client document."),
    ("\u2302", "Credit Unions and Banks",
     "Add authentication to every loan signing event."),
    ("\u2261", "Debt, Legal and Media",
     "Use RON for any legal notice that needs notarizing."),
    ("\u25c6", "Integration Partners",
     "Embed SIGNiX as the compliance-grade option in your platform."),
]

for (nx, ny), (icon, seg, behavior) in zip(NODE_CENTERS, segments):
    node_card(s, nx, ny, icon, seg, behavior)

# ── Save ─────────────────────────────────────────────────────────────────────
OUT = ("/Users/chris/Desktop/AI Summit 2026 Q2/"
       "proj-template-and-lease-SIGNiX-app/PROJECT-DOCS/"
       "SIGNiX_SlideC_Audiences.pptx")
prs.save(OUT)
print(f"Saved: {OUT}")
